#!/usr/bin/env python3
"""Ingest a local file into raw/.

Supported formats:
    .md / .markdown  - verbatim copy
    .txt             - wrapped in a code fence
    .html / .htm     - markdownify
    .pdf             - pypdf text extraction
    .docx            - python-docx paragraphs + tables
    .xlsx            - openpyxl sheets as tables
    .csv             - stdlib csv as a single table
    .pptx            - python-pptx slide text + tables
    .png/.jpg/.webp/.gif - the file IS the source (goes straight to raw/images/)

Usage:
    python3 fetch_local.py --wiki-root /path/to/wiki --path /path/to/source
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
import shutil
import sys
import zipfile
from pathlib import Path
from typing import Optional

from _deps import require

require(["markdownify", "bs4"])

from bs4 import BeautifulSoup
from markdownify import markdownify

from config import ConfigError, load_config
from raw_store import read_previous_source_metadata, write_raw_if_changed


_slug_re = re.compile(r"[^\w\-]+", re.UNICODE)
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}


def slugify(name: str) -> str:
    import unicodedata

    normalized = unicodedata.normalize("NFKD", name)
    ascii_only = normalized.encode("ascii", "ignore").decode("ascii")
    lowered = ascii_only.lower().strip()
    dashed = _slug_re.sub("-", lowered).strip("-")
    return dashed or "untitled"


def sanitize_cell(value) -> str:
    """Stringify a table cell (openpyxl values can be None/int/float/datetime)."""
    text = "" if value is None else str(value)
    return text.replace("\n", " ").replace("|", "\\|").strip()


def render_table(rows: list[list[str]]) -> list[str]:
    """Render rows of sanitized cell strings as a GitHub-flavored Markdown table.

    Rows are padded to the widest row's length so jagged input (e.g. CSV rows
    with missing trailing fields) doesn't misalign columns.
    """
    if not rows:
        return []
    width = max(len(r) for r in rows)
    padded = [r + [""] * (width - len(r)) for r in rows]
    col_widths = [max(len(c) for c in col) for col in zip(*padded)]
    header = "| " + " | ".join(padded[0]) + " |"
    sep = "|" + "|".join("-" * (w + 2) for w in col_widths) + "|"
    body_rows = ["| " + " | ".join(r) + " |" for r in padded[1:]]
    return [header, sep, *body_rows]


def parse_md(path: Path) -> tuple[str, list[Path]]:
    body = path.read_text(encoding="utf-8", errors="replace")
    if not body.startswith("#"):
        title = path.stem.replace("-", " ").replace("_", " ").title()
        body = f"# {title}\n\n" + body
    return body, []


def parse_txt(path: Path) -> tuple[str, list[Path]]:
    body = path.read_text(encoding="utf-8", errors="replace")
    title = path.stem.replace("-", " ").replace("_", " ").title()
    fenced = "```\n" + body.rstrip() + "\n```\n"
    return f"# {title}\n\n{fenced}", []


def parse_html(path: Path) -> tuple[str, list[str]]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    image_hints = [img.get("src", "") for img in soup.find_all("img") if img.get("src")]
    title = (soup.title.string if soup.title else path.stem) or path.stem
    title = title.strip()
    body_html = str(soup.body) if soup.body else str(soup)
    md = markdownify(body_html, heading_style="ATX", bullets="-").strip()
    body = f"# {title}\n\n{md}\n"
    return body, image_hints


def parse_pdf(path: Path, images_dir: Path) -> tuple[str, list[Path]]:
    require(["pypdf"])
    from pypdf import PdfReader
    from pypdf.errors import PyPdfError

    try:
        reader = PdfReader(str(path))
    except PyPdfError as e:
        raise SystemExit(f"ERROR: could not read PDF {path}: {e}")

    if reader.is_encrypted:
        try:
            if reader.decrypt("") == 0:
                raise SystemExit(
                    f"ERROR: PDF {path} is encrypted with a password. "
                    "Decrypt it first."
                )
        except NotImplementedError:
            raise SystemExit(
                f"ERROR: PDF {path} uses an unsupported encryption scheme."
            )

    parts: list[str] = []
    saved_images: list[Path] = []
    idx = 0

    for page_num, page in enumerate(reader.pages, start=1):
        text = ""
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            text = ""

        parts.append(f"## Page {page_num}\n\n{text.strip()}\n")

        try:
            page_images = list(page.images or [])
        except Exception:  # noqa: BLE001
            page_images = []

        for image_object in page_images:
            data = getattr(image_object, "data", None)
            name = getattr(image_object, "name", None) or f"page{page_num}-img"
            if data is None:
                continue
            ext = Path(name).suffix.lower() or ".png"
            if ext not in IMAGE_EXTS:
                ext = ".png"
            image_path = images_dir / f"{idx}{ext}"
            images_dir.mkdir(parents=True, exist_ok=True)
            image_path.write_bytes(data)
            saved_images.append(image_path)
            idx += 1

    title = path.stem.replace("-", " ").replace("_", " ").title()
    body = f"# {title}\n\n" + "\n---\n\n".join(parts)
    return body, saved_images


def parse_docx(path: Path, images_dir: Path) -> tuple[str, list[Path]]:
    require(["docx"])
    from docx import Document

    doc = Document(str(path))
    lines: list[str] = []

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        style = getattr(para.style, "name", "") or ""
        if not text:
            lines.append("")
            continue
        if style.startswith("Heading"):
            level_str = style.replace("Heading", "").strip() or "1"
            try:
                level = max(1, min(6, int(level_str)))
            except ValueError:
                level = 2
            lines.append("#" * level + " " + text)
        else:
            lines.append(text)
        lines.append("")

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)
        if not rows:
            continue
        widths = [max(len(c) for c in col) for col in zip(*rows)]
        header = "| " + " | ".join(rows[0]) + " |"
        sep = "|" + "|".join("-" * (w + 2) for w in widths) + "|"
        body_rows = ["| " + " | ".join(r) + " |" for r in rows[1:]]
        lines.append(header)
        lines.append(sep)
        lines.extend(body_rows)
        lines.append("")

    saved_images: list[Path] = []
    idx = 0
    with zipfile.ZipFile(path) as z:
        media_entries = sorted(n for n in z.namelist() if n.startswith("word/media/"))
        for name in media_entries:
            ext = Path(name).suffix.lower() or ".png"
            if ext not in IMAGE_EXTS:
                continue
            data = z.read(name)
            images_dir.mkdir(parents=True, exist_ok=True)
            image_path = images_dir / f"{idx}{ext}"
            image_path.write_bytes(data)
            saved_images.append(image_path)
            idx += 1

    title = path.stem.replace("-", " ").replace("_", " ").title()
    body = f"# {title}\n\n" + "\n".join(lines).strip() + "\n"
    return body, saved_images


def parse_xlsx(path: Path, images_dir: Path) -> tuple[str, list[Path]]:
    require(["openpyxl"])
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts: list[str] = []

    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        rows = [
            [sanitize_cell(c) for c in row] for row in ws.iter_rows(values_only=True)
        ]
        # Some sheets are blank section dividers — render the heading only.
        if any(cell for row in rows for cell in row):
            parts.append("")
            parts.extend(render_table(rows))
        parts.append("")

    saved_images: list[Path] = []
    idx = 0
    with zipfile.ZipFile(path) as z:
        media_entries = sorted(n for n in z.namelist() if n.startswith("xl/media/"))
        for name in media_entries:
            ext = Path(name).suffix.lower() or ".png"
            if ext not in IMAGE_EXTS:
                continue
            data = z.read(name)
            images_dir.mkdir(parents=True, exist_ok=True)
            image_path = images_dir / f"{idx}{ext}"
            image_path.write_bytes(data)
            saved_images.append(image_path)
            idx += 1

    title = path.stem.replace("-", " ").replace("_", " ").title()
    body = f"# {title}\n\n" + "\n".join(parts).strip() + "\n"
    return body, saved_images


def parse_csv(path: Path, images_dir: Path) -> tuple[str, list[Path]]:
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        sample = f.read(4096)
        f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        rows = [[sanitize_cell(c) for c in row] for row in csv.reader(f, dialect)]

    title = path.stem.replace("-", " ").replace("_", " ").title()
    table = render_table(rows)
    body = f"# {title}\n\n" + "\n".join(table).strip() + "\n"
    return body, []


def parse_pptx(path: Path, images_dir: Path) -> tuple[str, list[Path]]:
    require(["pptx"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    parts: list[str] = []
    saved_images: list[Path] = []
    idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = ""
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
        title_id = title_shape.shape_id if title_shape is not None else None

        heading = f"## Slide {slide_num}" + (f": {title_text}" if title_text else "")
        parts.append(heading)
        parts.append("")

        for shape in slide.shapes:
            if shape.shape_id == title_id:
                continue  # already used as the slide heading
            if shape.has_table:
                rows = [
                    [sanitize_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                parts.extend(render_table(rows))
                parts.append("")
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                for line in text.splitlines():
                    line = line.strip()
                    if line:
                        parts.append(f"- {line}")
                if text:
                    parts.append("")
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                except ValueError:
                    # Linked (not embedded) image, or a picture placeholder
                    # with no image part yet — nothing to extract.
                    continue
                ext = f".{image.ext}" if image.ext else ".png"
                if ext not in IMAGE_EXTS:
                    ext = ".png"
                images_dir.mkdir(parents=True, exist_ok=True)
                image_path = images_dir / f"{idx}{ext}"
                image_path.write_bytes(image.blob)
                saved_images.append(image_path)
                idx += 1

    title = path.stem.replace("-", " ").replace("_", " ").title()
    body = f"# {title}\n\n" + "\n".join(parts).strip() + "\n"
    return body, saved_images


def parse_image(
    path: Path, slug: str, raw_dir: Path
) -> tuple[str, list[Path]]:
    images_dir = raw_dir / "images" / slug
    images_dir.mkdir(parents=True, exist_ok=True)
    ext = path.suffix.lower()
    if ext not in IMAGE_EXTS:
        raise SystemExit(f"ERROR: unsupported image extension: {ext}")
    dest = images_dir / f"0{ext}"
    shutil.copyfile(path, dest)

    title = path.stem.replace("-", " ").replace("_", " ").title()
    md = (
        f"# {title}\n\n"
        f"![{title}](images/{slug}/{dest.name})\n\n"
        f"See `raw/images/{slug}/0.md` for the description.\n"
    )
    return md, [dest]


PARSERS = {
    ".md": ("copy", parse_md),
    ".markdown": ("copy", parse_md),
    ".txt": ("copy", parse_txt),
    ".html": ("html", parse_html),
    ".htm": ("html", parse_html),
    ".pdf": ("pdf", parse_pdf),
    ".docx": ("docx", parse_docx),
    ".xlsx": ("xlsx", parse_xlsx),
    ".csv": ("csv", parse_csv),
    ".pptx": ("pptx", parse_pptx),
}

# Parsers that take (path, images_dir) and may extract embedded images.
IMAGE_CAPABLE_KINDS = {"pdf", "docx", "xlsx", "pptx"}


def hash_source_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while chunk := f.read(65536):
            h.update(chunk)
    return h.hexdigest()


def read_previous_source_sha(raw_dir: Path, slug: str) -> Optional[str]:
    src_path = raw_dir / f"{slug}.source.json"
    if not src_path.exists():
        return None
    try:
        with src_path.open("r", encoding="utf-8") as f:
            return json.load(f).get("source_sha256")
    except (json.JSONDecodeError, OSError):
        return None


def resolve_slug_collision(
    raw_dir: Path, slug: str, original_path: str, explicit_slug: bool
) -> str:
    """Return a slug that isn't already owned by a different local file.

    The default slug is just the filename stem, so two different files that
    happen to share one (projectA/README.md, projectB/README.md) collide —
    write_raw_if_changed only ever compares bytes, never original_path, so the
    second ingest would silently overwrite the first's raw file. Unlike
    Confluence's page_id there's no natural short identifier here, so the
    disambiguator is a short hash of the resolved path.

    When the slug came from an explicit --slug, respect it — the user asked
    for that name deliberately — but still warn so the mismatch isn't silent.
    """
    prior = read_previous_source_metadata(raw_dir, slug)
    if not prior:
        return slug
    prior_path = str(prior.get("original_path") or "")
    if not prior_path or prior_path == original_path:
        return slug  # same file — keep the slug

    if explicit_slug:
        print(
            f"WARNING: --slug {slug!r} is already used by {prior_path} — a "
            f"different file. Proceeding will overwrite it with {original_path}.",
            file=sys.stderr,
        )
        return slug

    digest = hashlib.sha256(original_path.encode("utf-8")).hexdigest()[:8]
    candidate = f"{slug}-{digest}"
    prior2 = read_previous_source_metadata(raw_dir, candidate)
    if prior2:
        prior2_path = str(prior2.get("original_path") or "")
        if prior2_path and prior2_path != original_path:
            raise SystemExit(
                f"ERROR: slug collision could not be resolved for {original_path}: "
                f"both {slug!r} and {candidate!r} are already owned by other files "
                f"({prior_path}, {prior2_path}). Pass --slug to disambiguate manually."
            )
    print(
        f"WARNING: slug {slug!r} is already used by {prior_path} — a different "
        f"file with the same name. Using {candidate!r} for {original_path} instead. "
        "Pass --slug to choose your own.",
        file=sys.stderr,
    )
    return candidate


def copy_into_raw(source: Path, raw_dir: Path, slug: str, ext: str) -> Path:
    """Copy the original file into raw/<slug><ext> and return the copy's path.

    This makes the wiki self-contained: the diff check (source_sha256) and any
    re-parse read the copy, not the external original — so a moved or deleted
    original never breaks re-ingest, and a fresh clone owns its sources.

    Idempotent: if `source` already IS the copy (re-ingesting the in-raw file),
    the copy is skipped. `ext` is lowercased for a stable filename.
    """
    raw_dir.mkdir(parents=True, exist_ok=True)
    dest = raw_dir / f"{slug}{ext.lower()}"
    if source.resolve() == dest.resolve():
        return dest  # already the copy — nothing to do
    shutil.copyfile(source, dest)
    return dest


def main() -> int:
    parser = argparse.ArgumentParser(description="Ingest a local file into raw/")
    parser.add_argument("--wiki-root", type=Path, required=True)
    parser.add_argument("--path", type=Path, required=True)
    parser.add_argument("--slug", default=None, help="Override the slug (default: filename stem)")
    args = parser.parse_args()

    try:
        cfg = load_config(args.wiki_root)
    except ConfigError as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 1

    source = args.path.expanduser().resolve()
    if not source.exists():
        print(f"ERROR: file not found: {source}", file=sys.stderr)
        return 1

    ext = source.suffix.lower()
    slug = args.slug or slugify(source.stem)
    raw_dir = cfg.raw_dir
    raw_dir.mkdir(parents=True, exist_ok=True)
    original_filename = source.name
    original_path = str(source)
    slug = resolve_slug_collision(raw_dir, slug, original_path, bool(args.slug))
    images_dir = raw_dir / "images" / slug

    # Copy the original into raw/ so the wiki owns its source and the diff check
    # never depends on the external file. Two extensions are exempt from the
    # root copy:
    #   - Images: they have a dedicated home (raw/images/<slug>/) via parse_image.
    #   - Markdown: the rendered raw/<slug>.md IS the faithful in-wiki copy, so a
    #     separate raw/<slug>.md copy would collide with write_raw_if_changed's
    #     own output. (.txt renders to .md, so no collision — it's still copied.)
    is_image = ext in IMAGE_EXTS
    is_markdown = ext in {".md", ".markdown"}
    if is_image:
        parse_source = source
        rel_path = None  # set after parse_image copies it
    elif is_markdown:
        parse_source = source
        rel_path = f"{raw_dir.name}/{slug}.md"  # the rendered file is the copy
    else:
        parse_source = copy_into_raw(source, raw_dir, slug, ext)
        rel_path = f"{raw_dir.name}/{parse_source.name}"

    # Hash the copied-in file (or the image original) — the diff baseline is the
    # artifact the wiki owns, not the external path.
    source_sha = hash_source_file(parse_source)
    prior_sha = read_previous_source_sha(raw_dir, slug)

    if prior_sha == source_sha:
        # Fast-path: the source bytes haven't changed since last ingest.
        # Don't re-parse (PDF/DOCX parsing is expensive), don't rewrite anything.
        md_path = raw_dir / f"{slug}.md"
        src_path = raw_dir / f"{slug}.source.json"
        print(
            json.dumps(
                {
                    "slug": slug,
                    "title": slug.replace("-", " ").title(),
                    "raw_md": str(md_path),
                    "source_json": str(src_path),
                    "image_hints": [],
                    "saved_images": [],
                    "status": "unchanged",
                    "source_sha256": source_sha,
                },
                indent=2,
                ensure_ascii=False,
            )
        )
        return 0

    saved_images: list[Path] = []
    image_hints: list = []

    if is_image:
        markdown, saved_images = parse_image(parse_source, slug, raw_dir)
        kind = "image"
        # parse_image copied the bytes into raw/images/<slug>/0<ext>.
        rel_path = f"{raw_dir.name}/images/{slug}/{saved_images[0].name}"
    elif ext in PARSERS:
        kind, parser_fn = PARSERS[ext]
        if kind in IMAGE_CAPABLE_KINDS or kind == "csv":
            markdown, saved_images = parser_fn(parse_source, images_dir)
        elif kind == "html":
            markdown, image_hints = parser_fn(parse_source)
        else:
            markdown, _ = parser_fn(parse_source)
    else:
        # No parser for this extension (e.g. legacy binary .xls/.ppt/.doc). The
        # original is still copied into raw/ and versioned; Claude synthesizes
        # the .md from it. Emit a minimal placeholder body so the raw artifact
        # is self-describing.
        kind = ext.lstrip(".") or "file"
        title = source.stem.replace("-", " ").replace("_", " ").title()
        markdown = (
            f"# {title}\n\n"
            f"_Original file: `{rel_path}` "
            f"({original_filename}). No text extractor for `{ext}` — "
            f"synthesize wiki content directly from the source file._\n"
        )

    metadata = {
        "type": "local",
        "kind": kind,
        "path": rel_path,
        "filename": original_filename,
        "original_filename": original_filename,
        "original_path": original_path,
        "source_sha256": source_sha,
        "image_hints": image_hints,
    }

    result = write_raw_if_changed(raw_dir, slug, markdown.strip() + "\n", metadata)

    print(
        json.dumps(
            {
                "slug": slug,
                "title": slug.replace("-", " ").title(),
                "raw_md": result["raw_md"],
                "source_json": result["source_json"],
                "image_hints": image_hints,
                "saved_images": [str(p) for p in saved_images],
                "status": result["status"],
                "source_sha256": source_sha,
                "content_sha256": result["content_sha256"],
            },
            indent=2,
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
